import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
import seaborn as sns
import os
from typing import Dict, Any, List, Tuple
import logging
import numpy as np
import matplotlib.cm as cm

logger = logging.getLogger(__name__)

def validate_and_prepare_data(df: pd.DataFrame, x_col: str, y_col: str, chart_type: str) -> Tuple[pd.DataFrame, str, str]:
    """Validate and prepare data for visualization."""
    # Check for empty dataframe
    if df.empty:
        raise ValueError("The input file is empty")

    # Validate columns exist
    if x_col not in df.columns or y_col not in df.columns:
        raise ValueError(f"Columns {x_col} and/or {y_col} not found in the file")

    # Handle null values
    null_count = df[[x_col, y_col]].isna().sum()
    if null_count.any():
        logger.warning(f"Found null values: {null_count.to_dict()}")
        df = df.dropna(subset=[x_col, y_col])
        if df.empty:
            raise ValueError("No valid data points after removing null values")

    # Chart-specific validations and preparations
    if chart_type in ['bar', 'line', 'scatter']:
        if not pd.api.types.is_numeric_dtype(df[y_col]):
            try:
                df[y_col] = pd.to_numeric(df[y_col], errors='coerce')
                if df[y_col].isna().all():
                    raise ValueError(f"Y-axis column '{y_col}' must contain numeric data for {chart_type} chart")
                df = df.dropna(subset=[y_col])
            except:
                raise ValueError(f"Cannot convert Y-axis column '{y_col}' to numeric data")

    elif chart_type == 'pie':
        unique_values = df[x_col].nunique()
        if unique_values > 10:
            raise ValueError(f"Pie chart not suitable for {unique_values} categories. Maximum recommended is 10.")
        if not pd.api.types.is_numeric_dtype(df[y_col]):
            try:
                df[y_col] = pd.to_numeric(df[y_col], errors='coerce')
                df = df.dropna(subset=[y_col])
                if df.empty:
                    raise ValueError("No valid numeric data for pie chart")
            except:
                raise ValueError("Y-axis column must be numeric for pie charts")

    # Handle datetime data
    if pd.api.types.is_datetime64_any_dtype(df[x_col]):
        df[x_col] = df[x_col].dt.strftime('%Y-%m-%d')

    # Sample large datasets
    if len(df) > 1000:
        df = df.sample(n=1000, random_state=42)
        logger.info("Dataset sampled to 1000 points for better performance")

    return df, x_col, y_col

def generate_chart_data(file_path: str, x_col: str, y_col: str, chart_type: str) -> List[Dict[str, Any]]:
    """Generate chart data from the input file and selected columns."""
    try:
        # Read the file based on extension
        try:
            if file_path.endswith(('.xlsx', '.xls')):
                df = pd.read_excel(file_path)
            else:
                df = pd.read_csv(file_path)
        except Exception as e:
            raise ValueError(f"Failed to read file: {str(e)}")

        # Validate and prepare data
        df, x_col, y_col = validate_and_prepare_data(df, x_col, y_col, chart_type)

        # Generate chart data based on chart type
        if chart_type == "pie":
            chart_data = (
                df.groupby(x_col)[y_col]
                .agg('sum')
                .reset_index()
                .to_dict('records')
            )
        else:
            # For other charts, sort by x_col and handle numeric x-axis
            if pd.api.types.is_numeric_dtype(df[x_col]):
                df = df.sort_values(by=x_col)
            chart_data = df[[x_col, y_col]].to_dict('records')

        # Validate final data
        if not chart_data:
            raise ValueError("No valid data points generated for the chart")

        logger.info(f"Successfully generated {chart_type} chart data with {len(chart_data)} points")
        return chart_data

    except Exception as e:
        logger.error(f"Error generating chart data: {str(e)}")
        raise ValueError(f"Failed to generate chart data: {str(e)}")

def export_to_ppt(chart_data: List[Dict[str, Any]], chart_type: str, x_column: str, y_column: str) -> str:
    """Create a PowerPoint presentation with the chart."""
    try:
        # Get the absolute path to the uploads directory
        base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        uploads_dir = os.path.join(base_dir, "uploads")
        
        # Ensure uploads directory exists
        os.makedirs(uploads_dir, exist_ok=True)
        
        # Create a new presentation
        prs = Presentation()
        
        # Add a title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = f"{y_column} vs {x_column}"
        subtitle.text = f"Chart Type: {chart_type.capitalize()}"
        
        # Create a new slide for the chart
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        # Set up the plot with error handling
        plt.figure(figsize=(10, 6))
        try:
            # Convert data to x and y series
            x_data = [item[x_column] for item in chart_data]
            y_data = [item[y_column] for item in chart_data]
            
            # Set dark blue color
            dark_blue = '#1e40af'  # Tailwind dark blue-700
            
            if chart_type == "bar":
                plt.bar(x_data, y_data, color=dark_blue)
            elif chart_type == "line":
                plt.plot(x_data, y_data, color=dark_blue, linewidth=2)
            elif chart_type == "scatter":
                plt.scatter(x_data, y_data, color=dark_blue)
            elif chart_type == "pie":
                # Use a color palette for the pie slices
                colors = cm.get_cmap('tab20').colors
                num_slices = len(x_data)
                slice_colors = colors[:num_slices]
                plt.pie(y_data, labels=x_data, autopct='%1.1f%%', colors=slice_colors)
            
            plt.xlabel(x_column)
            plt.ylabel(y_column)
            plt.title(f"{y_column} vs {x_column}")
            
            # Handle long labels
            if chart_type != "pie":
                plt.xticks(rotation=45, ha='right')
            
            # Adjust layout to prevent label cutoff
            plt.tight_layout()
            
        except Exception as e:
            plt.close()
            raise ValueError(f"Failed to create chart: {str(e)}")
        
        # Save the chart as a temporary image
        temp_image = os.path.join(uploads_dir, "temp_chart.png")
        try:
            plt.savefig(temp_image, bbox_inches='tight', dpi=300)
        except Exception as e:
            plt.close()
            raise ValueError(f"Failed to save chart image: {str(e)}")
        finally:
            plt.close()
        
        # Add the chart to the slide
        try:
            left = Inches(1)
            top = Inches(1)
            width = Inches(8)
            height = Inches(5)
            slide.shapes.add_picture(temp_image, left, top, width, height)
        except Exception as e:
            if os.path.exists(temp_image):
                os.remove(temp_image)
            raise ValueError(f"Failed to add chart to slide: {str(e)}")
        
        # Save the presentation
        output_path = os.path.join(uploads_dir, "chart_presentation.pptx")
        try:
            prs.save(output_path)
        except Exception as e:
            if os.path.exists(temp_image):
                os.remove(temp_image)
            raise ValueError(f"Failed to save PowerPoint file: {str(e)}")
        
        # Clean up the temporary image
        if os.path.exists(temp_image):
            try:
                os.remove(temp_image)
            except Exception as e:
                logger.warning(f"Failed to remove temporary image: {str(e)}")
        
        return output_path
        
    except Exception as e:
        # Clean up temporary files in case of error
        temp_image = os.path.join(uploads_dir, "temp_chart.png")
        if os.path.exists(temp_image):
            try:
                os.remove(temp_image)
            except:
                pass
        logger.error(f"Error in PPT export: {str(e)}")
        raise ValueError(f"Failed to export to PPT: {str(e)}")